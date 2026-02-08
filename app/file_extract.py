from __future__ import annotations

import base64
import io
import os
import re
import shutil
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Iterable, List, Tuple
import uuid

import pytesseract
from PIL import Image, ImageFilter, ImageOps
from pypdf import PdfReader
import fitz  # pymupdf

from email import policy
from email.parser import BytesParser

_OCR_CACHE: dict = {}
_PASSPORT_DOC_REGEX = re.compile(
    r"((?:passport|document)\b[\s\S]{0,200}?(?:E|P[A-E]|RA|N|M)\d{7}(?!\d)|(?:E|P[A-E]|RA|N|M)\d{7}(?!\d)[\s\S]{0,200}?\b(?:passport|document))",
    re.IGNORECASE,
)
_PASSPORT_DOC_REGEX_STRICT = re.compile(
    r"(?:^|[\s,;:\(\)\[\]\"'])(?:(?:passport|document)\b[\s\S]{0,200}?(?:E|P[A-E]|RA|N|M)\d{7}(?!\d)|(?:E|P[A-E]|RA|N|M)\d{7}(?!\d)[\s\S]{0,200}?\b(?:passport|document))(?:$|[\s,;:\(\)\[\]\"']|\.\s|\.$)",
    re.IGNORECASE,
)


@dataclass
class ExtractionResult:
    filename: str
    ok: bool
    message: str
    status: str = "success"


MAX_EXTRACTED_TEXT_CHARS = int(os.getenv("MAX_EXTRACTED_TEXT_CHARS", "400000"))
MAX_FILE_TEXT_CHARS = int(os.getenv("MAX_FILE_TEXT_CHARS", "200000"))
OCR_PDF_DPI = int(os.getenv("OCR_PDF_DPI", "200"))
MAX_OCR_PDF_PAGES = int(os.getenv("MAX_OCR_PDF_PAGES", "50"))
MAX_ARCHIVE_FILES = int(os.getenv("MAX_ARCHIVE_FILES", "2000"))
MAX_ARCHIVE_BYTES = int(os.getenv("MAX_ARCHIVE_BYTES", str(200 * 1024 * 1024)))
MAX_IMAGE_PIXELS = int(os.getenv("MAX_IMAGE_PIXELS", str(6_000_000)))
MIN_OCR_DIMENSION = int(os.getenv("MIN_OCR_DIMENSION", "24"))
MAX_MESSAGE_ATTACHMENTS = int(os.getenv("MAX_MESSAGE_ATTACHMENTS", "100"))
MAX_ATTACHMENT_BYTES = int(os.getenv("MAX_ATTACHMENT_BYTES", str(50 * 1024 * 1024)))


SUPPORTED_EXTENSIONS = {
    ".docx",
    ".pptx",
    ".xlsx",
    ".xls",
    ".rtf",
    ".odt",
    ".pdf",
    ".eml",
    ".msg",
    ".jpg",
    ".jpeg",
    ".png",
    ".tiff",
    ".tif",
    ".gif",
    ".zip",
    ".7z",
    ".rar",
}


def _safe_join(base: Path, *paths: str) -> Path:
    candidate = (base.joinpath(*paths)).resolve()
    if not str(candidate).startswith(str(base.resolve())):
        raise ValueError("Path traversal detected")
    return candidate


def _limit_extraction(root: Path, max_files: int, max_bytes: int) -> None:
    total = 0
    count = 0
    for path in root.rglob("*"):
        if path.is_file():
            count += 1
            if count > max_files:
                raise ValueError("Too many files extracted")
            total += path.stat().st_size
            if total > max_bytes:
                raise ValueError("Extracted content too large")


def _extract_zip(archive_path: Path, dest: Path, max_files: int, max_bytes: int) -> None:
    import zipfile

    with zipfile.ZipFile(archive_path) as zf:
        infos = zf.infolist()
        if len(infos) > max_files:
            raise ValueError("Too many files in zip")
        total = sum(i.file_size for i in infos)
        if total > max_bytes:
            raise ValueError("Zip content too large")
        for info in infos:
            if info.is_dir():
                continue
            target = _safe_join(dest, info.filename)
            target.parent.mkdir(parents=True, exist_ok=True)
            with zf.open(info) as src, open(target, "wb") as out:
                shutil.copyfileobj(src, out)


def _extract_7z(archive_path: Path, dest: Path, max_files: int, max_bytes: int) -> None:
    import py7zr

    with py7zr.SevenZipFile(archive_path, mode="r") as zf:
        names = zf.getnames()
        if len(names) > max_files:
            raise ValueError("Too many files in 7z")
        zf.extractall(path=dest)
    _limit_extraction(dest, max_files, max_bytes)


def _extract_rar(archive_path: Path, dest: Path, max_files: int, max_bytes: int) -> None:
    import rarfile

    with rarfile.RarFile(archive_path) as rf:
        infos = rf.infolist()
        if len(infos) > max_files:
            raise ValueError("Too many files in rar")
        total = sum(i.file_size for i in infos)
        if total > max_bytes:
            raise ValueError("Rar content too large")
        for info in infos:
            if info.is_dir():
                continue
            target = _safe_join(dest, info.filename)
            target.parent.mkdir(parents=True, exist_ok=True)
            with rf.open(info) as src, open(target, "wb") as out:
                shutil.copyfileobj(src, out)


def _prep_variants(image: Image.Image) -> list[Image.Image]:
    gray = ImageOps.grayscale(image)
    scale = 2
    gray = gray.resize((gray.width * scale, gray.height * scale), Image.Resampling.LANCZOS)
    gray = gray.filter(ImageFilter.MedianFilter(size=3))
    auto = ImageOps.autocontrast(gray)
    bw_160 = auto.point(lambda x: 0 if x < 160 else 255, mode="1")
    bw_190 = auto.point(lambda x: 0 if x < 190 else 255, mode="1")
    return [gray, auto, bw_160, bw_190]


def _prepare_ocr_image(image: Image.Image) -> Image.Image | None:
    width, height = image.size
    if width < MIN_OCR_DIMENSION or height < MIN_OCR_DIMENSION:
        return None

    pixels = width * height
    if pixels <= MAX_IMAGE_PIXELS:
        return image

    scale = (MAX_IMAGE_PIXELS / float(pixels)) ** 0.5
    new_size = (max(int(width * scale), MIN_OCR_DIMENSION), max(int(height * scale), MIN_OCR_DIMENSION))
    return image.resize(new_size, Image.Resampling.LANCZOS)


def _text_quality(text: str) -> float:
    if not text:
        return 0.0
    stripped = text.strip("\n\r\t ")
    if not stripped:
        return 0.0

    lines = [line.strip() for line in stripped.splitlines() if line.strip()]
    tokens = [tok for tok in re.split(r"\s+", stripped) if tok]
    if not tokens:
        return 0.0

    cleaned_tokens = [re.sub(r"[^A-Za-z0-9<>]", "", tok) for tok in tokens]
    cleaned_tokens = [tok for tok in cleaned_tokens if tok]
    if not cleaned_tokens:
        return 0.0

    letters = sum(ch.isalpha() for ch in stripped)
    digits = sum(ch.isdigit() for ch in stripped)
    spaces = sum(ch.isspace() for ch in stripped)
    total = len(stripped)
    alnum_ratio = (letters + digits) / max(total - spaces, 1)

    single_ratio = sum(1 for tok in cleaned_tokens if len(tok) == 1) / max(
        len(cleaned_tokens), 1
    )
    long_ratio = sum(1 for tok in cleaned_tokens if len(tok) >= 3) / max(
        len(cleaned_tokens), 1
    )

    if lines:
        line_token_counts = [len([t for t in re.split(r"\s+", line) if t]) for line in lines]
        avg_line_tokens = sum(line_token_counts) / len(line_token_counts)
        line_density = min(avg_line_tokens / 4.0, 1.0)
        rich_line_ratio = sum(1 for count in line_token_counts if count >= 3) / len(
            line_token_counts
        )
    else:
        line_density = 0.0
        rich_line_ratio = 0.0

    # Weighted blend tuned to penalize noisy OCR with many single-character lines.
    score = (
        0.20 * alnum_ratio
        + 0.25 * line_density
        + 0.25 * long_ratio
        + 0.20 * rich_line_ratio
        + 0.10 * (1.0 - single_ratio)
    )
    return max(0.0, min(score, 1.0))


def _score_ocr_text(text: str) -> tuple[float, str]:
    normalized = _normalize_pdf_text(text)
    return _text_quality(normalized), normalized


def _ocr_image(image: Image.Image) -> str:
    prepared = _prepare_ocr_image(image)
    if prepared is None:
        return ""

    if os.getenv("USE_EASYOCR", "0").lower() in {"1", "true", "yes", "on"}:
        try:
            import easyocr
            import numpy as np

            if "_easyocr_reader" not in _OCR_CACHE:
                _OCR_CACHE["_easyocr_reader"] = easyocr.Reader(["en"], gpu=False)
            reader = _OCR_CACHE["_easyocr_reader"]
            arr = np.array(prepared.convert("RGB"))
            results = reader.readtext(arr, detail=0, paragraph=True)
            return "\n".join([r for r in results if r]).strip()
        except Exception:
            pass

    lang = os.getenv("OCR_LANG", "eng")
    configs = [os.getenv("OCR_CONFIG", "--oem 1 --psm 6"), "--oem 1 --psm 4"]

    def _ocr_single_orientation(img: Image.Image) -> tuple[float, str]:
        best_text = ""
        best_quality = 0.0
        for variant in _prep_variants(img):
            for cfg in configs:
                text = pytesseract.image_to_string(variant, lang=lang, config=cfg)
                quality, normalized = _score_ocr_text(text)
                if quality > best_quality:
                    best_quality = quality
                    best_text = normalized
        return best_quality, best_text

    base_quality, base_text = _ocr_single_orientation(prepared)
    best_quality, best_text = base_quality, base_text

    # Retry rotated orientations only when initial OCR looks weak.
    if best_quality < 0.62:
        for angle in (90, 270):
            rotated = prepared.rotate(angle, expand=True)
            quality, text = _ocr_single_orientation(rotated)
            if quality > best_quality:
                best_quality = quality
                best_text = text

    # MRZ pass (bottom 35%) using OCRB if available
    mrz_text = ""
    try:
        ocrb_config = os.getenv(
            "OCR_MRZ_CONFIG",
            "--oem 1 --psm 6 -c tessedit_char_whitelist=ABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789<",
        )
        mrz_lang = os.getenv("OCR_MRZ_LANG", "ocrb")
        bottom = prepared.crop((0, int(prepared.height * 0.65), prepared.width, prepared.height))
        mrz_variant = _prep_variants(bottom)[2]
        mrz_text = pytesseract.image_to_string(mrz_variant, lang=mrz_lang, config=ocrb_config)
    except Exception:
        pass

    combined = "\n".join([t for t in [best_text.strip(), mrz_text.strip()] if t]).strip()
    return _normalize_pdf_text(combined)


def _normalize_pdf_text(text: str) -> str:
    if not text:
        return ""

    cleaned_lines: List[str] = []
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            cleaned_lines.append("")
            continue

        # Normalize whitespace while preserving line boundaries.
        line = re.sub(r"\s+", " ", line)

        # Recover common missing boundaries such as "BSBNumber", "2024to04/..."
        line = re.sub(r"([a-z])([A-Z])", r"\1 \2", line)
        line = re.sub(r"([A-Z]{2,})([A-Z][a-z])", r"\1 \2", line)
        line = re.sub(r"(\d)([A-Za-z])", r"\1 \2", line)
        line = re.sub(r"([A-Za-z])(\d)", r"\1 \2", line)

        # Keep punctuation spacing readable.
        line = re.sub(r"\s+([,.;:!?])", r"\1", line)
        line = re.sub(r"(?<=[A-Za-z])\.(?=[A-Za-z])", ". ", line)
        line = re.sub(r"\s+\)", ")", line)
        line = re.sub(r"\(\s+", "(", line)

        cleaned_lines.append(line.strip())

    text = "\n".join(cleaned_lines)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    text = _normalize_document_identifiers(text)
    text = _augment_passport_extraction(text)
    return text


def _normalize_document_identifiers(text: str) -> str:
    def _compact_digits(value: str) -> str:
        return re.sub(r"\D", "", value)

    # Normalize OCR variants such as "PA 4269170" -> "PA4269170"
    text = re.sub(
        r"\b((?:E|P[A-E]|RA|N|M))\s*[-–—]?\s*((?:\d[\s-]*){7})\b",
        lambda m: f"{m.group(1).upper()}{_compact_digits(m.group(2))}",
        text,
        flags=re.IGNORECASE,
    )
    # Normalize OCR variants such as "P A 4269170" -> "PA4269170"
    text = re.sub(
        r"\bP\s*([A-E])\s*((?:\d[\s-]*){7})\b",
        lambda m: f"P{m.group(1).upper()}{_compact_digits(m.group(2))}",
        text,
        flags=re.IGNORECASE,
    )
    # Normalize OCR variants such as "PAK 2691708" -> "PA2691708"
    text = re.sub(
        r"\b(P[A-E])[A-Z]\s*((?:\d[\s-]*){7})\b",
        lambda m: f"{m.group(1).upper()}{_compact_digits(m.group(2))}",
        text,
        flags=re.IGNORECASE,
    )
    return text


def _extract_passport_candidates(text: str) -> list[str]:
    candidates: list[str] = []
    seen: set[str] = set()

    for match in re.finditer(r"\b(?:E|P[A-E]|RA|N|M)\d{7}\b", text, flags=re.IGNORECASE):
        value = match.group(0).upper()
        if value not in seen:
            seen.add(value)
            candidates.append(value)

    for raw_line in text.splitlines():
        line = re.sub(r"[^A-Z0-9<]", "", raw_line.upper())
        if len(line) < 25 or "<" not in line:
            continue
        mrz = re.search(r"([A-Z0-9<]{9})\d[A-Z]{3}\d{6}", line)
        if not mrz:
            continue
        document_no = mrz.group(1).replace("<", "")
        if re.fullmatch(r"(?:E|P[A-E]|RA|N|M)\d{7}", document_no):
            if document_no not in seen:
                seen.add(document_no)
                candidates.append(document_no)

    return candidates


def _augment_passport_extraction(text: str) -> str:
    if not text:
        return text
    if _PASSPORT_DOC_REGEX.search(text) and _PASSPORT_DOC_REGEX_STRICT.search(text):
        return text

    candidates = _extract_passport_candidates(text)
    if not candidates:
        return text

    passport_signal = bool(
        re.search(r"(?i)\b(passport|document|australia|p<aus)\b", text)
        or re.search(r"(?i)\bobs?ervations\b", text)
    )
    if not passport_signal:
        return text

    candidate = candidates[0]
    hint = f"PASSPORT\nDOCUMENT NO\n{candidate}\n"
    combined = f"{text}\n{hint}"

    if _PASSPORT_DOC_REGEX.search(combined):
        return combined
    return text


def _extract_pdf_text_pymupdf(doc: fitz.Document) -> str:
    pages: List[str] = []
    for page in doc:
        page_text = page.get_text("text", sort=True) or ""
        page_text = _normalize_pdf_text(page_text)
        if page_text:
            pages.append(page_text)
    return "\n\n".join(pages).strip()


def _extract_pdf_text_pypdf(path: Path) -> str:
    text_parts: List[str] = []
    reader = PdfReader(str(path))
    for page in reader.pages:
        page_text = page.extract_text() or ""
        text_parts.append(page_text)
    return _normalize_pdf_text("\n".join(text_parts))


def _extract_pdf(path: Path) -> str:
    with fitz.open(str(path)) as doc:
        # Prefer PyMuPDF text extraction since it keeps page layout and spacing better
        # on bank statement style PDFs compared with generic stream extraction.
        fitz_text = _extract_pdf_text_pymupdf(doc)
        if fitz_text and _text_quality(fitz_text) > 0.62:
            return fitz_text

        # Fallback to pypdf stream extraction for text PDFs where fitz extraction is weak.
        pypdf_text = _extract_pdf_text_pypdf(path)
        if pypdf_text and _text_quality(pypdf_text) > 0.62:
            return pypdf_text

        # OCR fallback for scanned PDFs (bounded page count and DPI for memory safety)
        ocr_text_parts: List[str] = []
        for page_idx, page in enumerate(doc):
            if page_idx >= MAX_OCR_PDF_PAGES:
                ocr_text_parts.append(
                    f"[OCR truncated after {MAX_OCR_PDF_PAGES} pages]"
                )
                break
            page_best_text = ""
            page_best_quality = 0.0

            try:
                textpage = page.get_textpage_ocr(language=os.getenv("OCR_LANG", "eng"))
                page_text_raw = textpage.extractText().strip()
                if page_text_raw:
                    tp_quality, tp_text = _score_ocr_text(page_text_raw)
                    page_best_quality = tp_quality
                    page_best_text = tp_text
            except Exception:
                pass

            pix = page.get_pixmap(dpi=OCR_PDF_DPI)
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            try:
                image_text = _ocr_image(img)
                iq, inorm = _score_ocr_text(image_text)
                if iq > page_best_quality:
                    page_best_quality = iq
                    page_best_text = inorm
            finally:
                img.close()

            if page_best_text:
                ocr_text_parts.append(page_best_text)

        if ocr_text_parts:
            return "\n".join([p for p in ocr_text_parts if p]).strip()

    return ""


def _extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text:
                    parts.append(shape.text)
    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: List[str] = []
    for sheet in wb.worksheets:
        parts.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell) for cell in row if cell is not None]
            if values:
                parts.append("\t".join(values))
    return "\n".join(parts)


def _extract_xls(path: Path) -> str:
    import xlrd

    wb = xlrd.open_workbook(str(path))
    parts: List[str] = []
    for sheet in wb.sheets():
        parts.append(f"# Sheet: {sheet.name}")
        for row_idx in range(sheet.nrows):
            row = sheet.row_values(row_idx)
            values = [str(cell) for cell in row if cell not in ("", None)]
            if values:
                parts.append("\t".join(values))
    return "\n".join(parts)


def _extract_rtf(path: Path) -> str:
    from striprtf.striprtf import rtf_to_text

    return rtf_to_text(path.read_text(encoding="utf-8", errors="ignore"))


def _extract_odt(path: Path) -> str:
    import odf.text
    import odf.opendocument

    doc = odf.opendocument.load(str(path))
    parts: List[str] = []
    for paragraph in doc.getElementsByType(odf.text.P):
        parts.append(str(paragraph))
    return "\n".join(parts)


def _extract_doc_legacy(path: Path) -> str:
    try:
        import textract
    except Exception as exc:  # pragma: no cover - optional
        raise RuntimeError("textract is required for .doc/.ppt files") from exc
    data = textract.process(str(path))
    return data.decode("utf-8", errors="replace")


def _extract_image(path: Path) -> str:
    image = Image.open(path)
    try:
        return _ocr_image(image)
    finally:
        image.close()


def _emit_progress(
    progress_callback: Callable[[ExtractionResult], None] | None,
    filename: str,
    ok: bool,
    message: str,
    status: str,
) -> ExtractionResult:
    report = ExtractionResult(filename=filename, ok=ok, message=message, status=status)
    if progress_callback:
        progress_callback(report)
    return report


def _truncate_text(text: str, limit: int) -> tuple[str, bool]:
    if len(text) <= limit:
        return text, False
    truncated = text[:limit]
    return truncated + f"\n\n[Truncated to {limit} characters]", True


def _extract_archive(
    path: Path,
    temp_dir: Path,
    progress_callback: Callable[[ExtractionResult], None] | None = None,
    display_prefix: str | None = None,
) -> Tuple[str, List[ExtractionResult]]:
    extract_dir = temp_dir / "archive"
    extract_dir.mkdir(parents=True, exist_ok=True)

    if path.suffix.lower() == ".zip":
        _extract_zip(path, extract_dir, MAX_ARCHIVE_FILES, MAX_ARCHIVE_BYTES)
    elif path.suffix.lower() == ".7z":
        _extract_7z(path, extract_dir, MAX_ARCHIVE_FILES, MAX_ARCHIVE_BYTES)
    elif path.suffix.lower() == ".rar":
        _extract_rar(path, extract_dir, MAX_ARCHIVE_FILES, MAX_ARCHIVE_BYTES)
    else:
        raise ValueError("Unsupported archive")

    files_to_process = [f for f in sorted(extract_dir.rglob("*")) if f.is_file()]
    root_name = display_prefix or path.name
    display_names = {file: f"{root_name}/{file.relative_to(extract_dir).as_posix()}" for file in files_to_process}

    reports_by_name: dict[str, ExtractionResult] = {}
    for file in files_to_process:
        display_name = display_names[file]
        queued = _emit_progress(progress_callback, display_name, True, "Queued", "info")
        reports_by_name[display_name] = queued

    parts: List[str] = []
    for file in files_to_process:
        display_name = display_names[file]
        try:
            suffix = file.suffix.lower()
            if suffix in {".zip", ".7z", ".rar", ".eml", ".msg"}:
                text, _ = _extract_file_with_reports(file, temp_dir, display_name, progress_callback)
            else:
                text = _extract_file_basic(file, temp_dir)
            text, was_truncated = _truncate_text(text, MAX_FILE_TEXT_CHARS)
            if text.strip():
                if text.lstrip().startswith("===== FILE:"):
                    parts.append(f"{text}\n")
                else:
                    parts.append(f"===== FILE: {display_name} =====\n{text}\n")
            message = "Extracted"
            if was_truncated:
                message = f"Extracted (truncated to {MAX_FILE_TEXT_CHARS} chars)"
            report = _emit_progress(progress_callback, display_name, True, message, "success")
            reports_by_name[display_name] = report
        except Exception as exc:
            report = _emit_progress(progress_callback, display_name, False, str(exc), "error")
            reports_by_name[display_name] = report
    return "\n".join(parts).strip(), list(reports_by_name.values())


def _extract_file_basic(path: Path, temp_dir: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix == ".docx":
        return _extract_docx(path)
    if suffix == ".doc":
        raise ValueError("Legacy .doc files are not supported in this build")
    if suffix == ".pptx":
        return _extract_pptx(path)
    if suffix == ".ppt":
        raise ValueError("Legacy .ppt files are not supported in this build")
    if suffix == ".xlsx":
        return _extract_xlsx(path)
    if suffix == ".xls":
        return _extract_xls(path)
    if suffix == ".rtf":
        return _extract_rtf(path)
    if suffix == ".odt":
        return _extract_odt(path)
    if suffix in {".jpg", ".jpeg", ".png", ".tif", ".tiff", ".gif"}:
        return _extract_image(path)

    raise ValueError(f"Unsupported file type: {suffix}")


def _extract_eml_itemized(
    path: Path,
    temp_dir: Path,
    display_name: str,
    progress_callback: Callable[[ExtractionResult], None] | None,
) -> Tuple[str, List[ExtractionResult]]:
    reports: List[ExtractionResult] = []
    parts: List[str] = []
    msg = BytesParser(policy=policy.default).parsebytes(path.read_bytes())

    body_name = f"{display_name}/body.txt"
    reports.append(_emit_progress(progress_callback, body_name, True, "Queued", "info"))

    headers = []
    if msg.get("subject"):
        headers.append(f"Subject: {msg.get('subject')}")
    if msg.get("from"):
        headers.append(f"From: {msg.get('from')}")
    if msg.get("to"):
        headers.append(f"To: {msg.get('to')}")

    body_parts: List[str] = []
    attachments: List[tuple[str, bytes]] = []

    if msg.is_multipart():
        for part in msg.walk():
            disp = part.get_content_disposition()
            ctype = (part.get_content_type() or "").lower()
            if disp == "attachment":
                filename = Path(part.get_filename() or "attachment").name
                payload = part.get_payload(decode=True) or b""
                attachments.append((filename, payload))
            elif ctype.startswith("text/"):
                text = part.get_content()
                if text:
                    body_parts.append(text)
            elif ctype.startswith("image/"):
                filename = Path(part.get_filename() or "inline-image").name
                payload = part.get_payload(decode=True) or b""
                attachments.append((filename, payload))
    else:
        text = msg.get_content()
        if text:
            body_parts.append(text)

    body_text = "\n".join(headers + body_parts).strip()
    body_text, body_truncated = _truncate_text(body_text, MAX_FILE_TEXT_CHARS)
    if body_text:
        parts.append(f"===== FILE: {body_name} =====\n{body_text}\n")
    body_msg = "Extracted"
    if body_truncated:
        body_msg = f"Extracted (truncated to {MAX_FILE_TEXT_CHARS} chars)"
    reports.append(_emit_progress(progress_callback, body_name, True, body_msg, "success"))

    attach_dir = temp_dir / f"message_attachments_{uuid.uuid4().hex}"
    attach_dir.mkdir(parents=True, exist_ok=True)
    for idx, (filename, payload) in enumerate(attachments):
        if idx >= MAX_MESSAGE_ATTACHMENTS:
            reports.append(
                _emit_progress(
                    progress_callback,
                    f"{display_name}/[attachments]",
                    False,
                    f"Attachment limit reached ({MAX_MESSAGE_ATTACHMENTS})",
                    "error",
                )
            )
            break
        if len(payload) > MAX_ATTACHMENT_BYTES:
            reports.append(
                _emit_progress(
                    progress_callback,
                    f"{display_name}/{filename}",
                    False,
                    f"Attachment too large (> {MAX_ATTACHMENT_BYTES} bytes)",
                    "error",
                )
            )
            continue
        attachment_path = attach_dir / filename
        attachment_path.parent.mkdir(parents=True, exist_ok=True)
        attachment_path.write_bytes(payload)
        attachment_display = f"{display_name}/{filename}"
        att_text, att_reports = _extract_file_with_reports(
            attachment_path, temp_dir, attachment_display, progress_callback
        )
        if att_text:
            parts.append(att_text)
        reports.extend(att_reports)

    return "\n".join([p for p in parts if p]).strip(), reports


def _extract_msg_itemized(
    path: Path,
    temp_dir: Path,
    display_name: str,
    progress_callback: Callable[[ExtractionResult], None] | None,
) -> Tuple[str, List[ExtractionResult]]:
    import extract_msg

    reports: List[ExtractionResult] = []
    parts: List[str] = []
    msg = extract_msg.Message(str(path))
    msg.process()

    body_name = f"{display_name}/body.txt"
    reports.append(_emit_progress(progress_callback, body_name, True, "Queued", "info"))

    body = []
    if msg.subject:
        body.append(f"Subject: {msg.subject}")
    if msg.sender:
        body.append(f"From: {msg.sender}")
    if msg.to:
        body.append(f"To: {msg.to}")
    if msg.body:
        body.append(msg.body)

    body_text = "\n".join(body).strip()
    body_text, body_truncated = _truncate_text(body_text, MAX_FILE_TEXT_CHARS)
    if body_text:
        parts.append(f"===== FILE: {body_name} =====\n{body_text}\n")
    body_msg = "Extracted"
    if body_truncated:
        body_msg = f"Extracted (truncated to {MAX_FILE_TEXT_CHARS} chars)"
    reports.append(_emit_progress(progress_callback, body_name, True, body_msg, "success"))

    attach_dir = temp_dir / f"msg_attachments_{uuid.uuid4().hex}"
    attach_dir.mkdir(parents=True, exist_ok=True)
    for idx, attachment in enumerate(msg.attachments):
        if idx >= MAX_MESSAGE_ATTACHMENTS:
            reports.append(
                _emit_progress(
                    progress_callback,
                    f"{display_name}/[attachments]",
                    False,
                    f"Attachment limit reached ({MAX_MESSAGE_ATTACHMENTS})",
                    "error",
                )
            )
            break

        filename = Path(attachment.longFilename or attachment.shortFilename or f"attachment_{idx}").name
        attachment.save(customPath=str(attach_dir))
        attachment_path = attach_dir / filename
        if not attachment_path.exists():
            reports.append(
                _emit_progress(
                    progress_callback,
                    f"{display_name}/{filename}",
                    False,
                    "Attachment save failed",
                    "error",
                )
            )
            continue
        if attachment_path.stat().st_size > MAX_ATTACHMENT_BYTES:
            reports.append(
                _emit_progress(
                    progress_callback,
                    f"{display_name}/{filename}",
                    False,
                    f"Attachment too large (> {MAX_ATTACHMENT_BYTES} bytes)",
                    "error",
                )
            )
            continue

        attachment_display = f"{display_name}/{filename}"
        att_text, att_reports = _extract_file_with_reports(
            attachment_path, temp_dir, attachment_display, progress_callback
        )
        if att_text:
            parts.append(att_text)
        reports.extend(att_reports)

    return "\n".join([p for p in parts if p]).strip(), reports


def _extract_file_with_reports(
    path: Path,
    temp_dir: Path,
    display_name: str,
    progress_callback: Callable[[ExtractionResult], None] | None,
) -> Tuple[str, List[ExtractionResult]]:
    suffix = path.suffix.lower()

    if suffix in {".zip", ".7z", ".rar"}:
        return _extract_archive(path, temp_dir, progress_callback=progress_callback, display_prefix=display_name)
    if suffix == ".eml":
        return _extract_eml_itemized(path, temp_dir, display_name, progress_callback)
    if suffix == ".msg":
        return _extract_msg_itemized(path, temp_dir, display_name, progress_callback)

    reports = [_emit_progress(progress_callback, display_name, True, "Queued", "info")]
    text = _extract_file_basic(path, temp_dir)
    text, was_truncated = _truncate_text(text, MAX_FILE_TEXT_CHARS)
    msg = "Extracted"
    if was_truncated:
        msg = f"Extracted (truncated to {MAX_FILE_TEXT_CHARS} chars)"
    reports.append(_emit_progress(progress_callback, display_name, True, msg, "success"))

    if text.strip():
        return f"===== FILE: {display_name} =====\n{text}\n", reports
    return "", reports


def extract_text_from_uploads(
    uploaded_files: Iterable,
    progress_callback: Callable[[ExtractionResult], None] | None = None,
) -> Tuple[str, List[ExtractionResult]]:
    combined_parts: List[str] = []
    results: List[ExtractionResult] = []
    total_chars = 0

    for uploaded in uploaded_files:
        filename = Path(uploaded.name).name
        suffix = Path(filename).suffix.lower()
        if suffix not in SUPPORTED_EXTENSIONS:
            report = ExtractionResult(
                filename=filename,
                ok=False,
                message="Unsupported file type",
                status="error",
            )
            results.append(report)
            if progress_callback:
                progress_callback(report)
            continue

        with tempfile.TemporaryDirectory() as tmpdir:
            temp_dir = Path(tmpdir)
            file_path = temp_dir / filename
            with open(file_path, "wb") as out:
                out.write(uploaded.getbuffer())
            try:
                text, file_reports = _extract_file_with_reports(
                    file_path,
                    temp_dir,
                    filename,
                    progress_callback,
                )
                remaining = MAX_EXTRACTED_TEXT_CHARS - total_chars
                if remaining <= 0:
                    results.extend(file_reports)
                    results.append(
                        _emit_progress(
                            progress_callback,
                            filename,
                            False,
                            f"Skipped adding text: total extracted text limit ({MAX_EXTRACTED_TEXT_CHARS}) reached",
                            "error",
                        )
                    )
                    continue
                text, was_total_truncated = _truncate_text(text, remaining)
                if text:
                    combined_parts.append(text)
                    total_chars += len(text)
                results.extend(file_reports)
                if was_total_truncated:
                    results.append(
                        _emit_progress(
                            progress_callback,
                            filename,
                            False,
                            f"Text truncated by total limit ({MAX_EXTRACTED_TEXT_CHARS})",
                            "error",
                        )
                    )
            except Exception as exc:
                results.append(
                    _emit_progress(progress_callback, filename, False, str(exc), "error")
                )

    return "\n".join(combined_parts).strip(), results
